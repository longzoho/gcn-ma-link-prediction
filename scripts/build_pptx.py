#!/usr/bin/env python3
"""
build_pptx.py — Generate a full 28-slide PPTX from the thesis defense outline.

Usage:
    .venv/bin/python scripts/build_pptx.py
    .venv/bin/python scripts/build_pptx.py --validate
"""
from __future__ import annotations

import argparse
import re
import sys
import unicodedata
from pathlib import Path
from typing import NamedTuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
REPO_ROOT = Path(__file__).resolve().parent.parent
OUTLINE_PATH = REPO_ROOT / "docs" / "slides" / "thesis_defense_outline.md"
OUTPUT_PATH = REPO_ROOT / "docs" / "slides" / "thesis_defense.pptx"
PLOTS_DIR = REPO_ROOT / "results" / "report" / "plots"

# ---------------------------------------------------------------------------
# Slide dimensions (16:9)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Color palette (academic / clean)
# ---------------------------------------------------------------------------
C_DARK_BLUE = RGBColor(0x1F, 0x35, 0x64)   # deep navy — titles
C_ACCENT = RGBColor(0x26, 0x72, 0xB6)       # mid-blue — accent
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
C_BODY_TEXT = RGBColor(0x26, 0x26, 0x26)    # near-black

# ---------------------------------------------------------------------------
# Known image mappings: slide-id (str) → list[filename]
# ---------------------------------------------------------------------------
SLIDE_IMAGES: dict[str, list[str]] = {
    "1":  ["dataset_snapshots_grid.png"],
    "5":  ["dataset_snapshots_grid.png"],
    "6":  ["edge_growth_density.png"],
    "7":  ["topology_map_2d.png"],
    "15": ["auc_comparison.png", "ap_comparison.png"],
    "16": ["ranking_heatmap.png"],
    "17": ["learning_curves_collegemsg.png", "learning_curves_mooc_actions.png"],
    "18": ["topology_map_2d_with_winners.png"],
    "19": ["beta_sensitivity.png"],
    "20": ["runtime_comparison.png"],
    "22": ["dataset_snapshots_grid.png"],
    "A2": [
        "learning_curves_bitcoinotc.png",
        "learning_curves_eut.png",
        "learning_curves_lastfm.png",
        "learning_curves_wikipedia.png",
    ],
}

# ---------------------------------------------------------------------------
# Latex / markdown → Unicode normalization
# ---------------------------------------------------------------------------
LATEX_MAP: list[tuple[str, str]] = [
    # multi-char replacements first (order matters)
    (r"\bar{\rho}", "ρ̄"),
    (r"\hat{S}", "Ŝ"),
    (r"\mathrm{LSTM}", "LSTM"),
    (r"\mathrm{GRU}", "GRU"),
    (r"\dots", "…"),
    (r"\cdot", "·"),
    (r"\cap", "∩"),
    (r"\deg", "deg"),
    (r"\sum", "Σ"),
    (r"\prod", "Π"),
    (r"\beta", "β"),
    (r"\rho", "ρ"),
]


def normalize_text(text: str) -> str:
    """Apply Unicode substitutions for inline LaTeX and clean markdown."""
    for latex, uni in LATEX_MAP:
        text = text.replace(latex, uni)
    # Strip remaining math delimiters
    text = text.replace("$", "")
    # Strip bold markers, keep text
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    # NFC normalize
    text = unicodedata.normalize("NFC", text)
    return text


# ---------------------------------------------------------------------------
# Data model
# ---------------------------------------------------------------------------
class SlideSpec(NamedTuple):
    slide_id: str    # "1".."22", "A1".."A6"
    title: str
    bullets: list[str]       # text bullets (IMAGE lines already removed)
    images: list[Path]       # resolved image paths to embed


# ---------------------------------------------------------------------------
# Parser
# ---------------------------------------------------------------------------
_HEADING_RE = re.compile(r"^##\s+Slide\s+(.+?)(?:\s+—\s+(.+))?$")
_IMAGE_RE = re.compile(r"^-\s+(?:IMAGE:\s*|(?:\[Khung trống cho\b[^\]]*\]\s*))(.+\.png)", re.IGNORECASE)
_BULLET_RE = re.compile(r"^-\s+(.+)$")


def parse_outline(path: Path) -> list[SlideSpec]:
    """Parse the markdown outline into SlideSpec objects."""
    slides: list[SlideSpec] = []
    current_id: str | None = None
    current_title: str = ""
    current_bullets: list[str] = []
    current_images: list[Path] = []

    def flush() -> None:
        if current_id is None:
            return
        # Resolve images from SLIDE_IMAGES registry (overrides inline IMAGE lines)
        # Use registry if present, else fall back to inline-parsed images.
        reg_images = SLIDE_IMAGES.get(current_id, [])
        resolved: list[Path] = []
        if reg_images:
            for fname in reg_images:
                p = PLOTS_DIR / fname
                if p.exists():
                    resolved.append(p)
        else:
            resolved = current_images

        slides.append(SlideSpec(
            slide_id=current_id,
            title=normalize_text(current_title),
            bullets=[normalize_text(b) for b in current_bullets],
            images=resolved,
        ))

    for raw_line in path.read_text(encoding="utf-8").splitlines():
        line = raw_line.rstrip()

        # Section dividers / blank
        if not line or line.startswith("---"):
            continue

        m = _HEADING_RE.match(line)
        if m:
            flush()
            raw_id_part = m.group(1).strip()
            raw_title_part = m.group(2).strip() if m.group(2) else raw_id_part
            # Normalize ID: "A1", "1", etc.
            current_id = raw_id_part.split()[0]  # "A1", "22", etc.
            current_title = raw_title_part
            current_bullets = []
            current_images = []
            continue

        if current_id is None:
            continue  # preamble text before first slide

        # Check for IMAGE lines (inline)
        img_m = _IMAGE_RE.match(line)
        if img_m:
            img_path_str = img_m.group(1).strip()
            # Could be a full relative path like results/report/plots/foo.png
            # or just a filename
            p = (REPO_ROOT / img_path_str) if "/" in img_path_str else (PLOTS_DIR / img_path_str)
            p = p.resolve()
            if p.exists():
                current_images.append(p)
            # Don't add to bullets
            continue

        # Regular bullets
        bul_m = _BULLET_RE.match(line)
        if bul_m:
            text = bul_m.group(1).strip()
            # Drop empty placeholder lines like "[Khung trống cho ...]"
            if re.match(r"\[Khung trống", text):
                continue
            current_bullets.append(text)

    flush()
    return slides


# ---------------------------------------------------------------------------
# PPTX helpers
# ---------------------------------------------------------------------------

def _set_slide_bg(slide, color: RGBColor) -> None:
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, text: str, left, top, width, height,
                  font_size: int = 28,
                  bold: bool = False,
                  color: RGBColor = C_BODY_TEXT,
                  align: PP_ALIGN = PP_ALIGN.LEFT,
                  wrap: bool = True) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _add_title_box(slide, title: str) -> None:
    """Standard title bar: left strip + title text."""
    # Accent bar
    left_bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        Inches(0.15), SLIDE_H,
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = C_ACCENT
    left_bar.line.fill.background()

    # Title background strip
    title_bg = slide.shapes.add_shape(
        1,
        Inches(0.15), Inches(0),
        SLIDE_W - Inches(0.15), Inches(1.05),
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = C_DARK_BLUE
    title_bg.line.fill.background()

    # Title text
    _add_text_box(
        slide, title,
        left=Inches(0.3), top=Inches(0.08),
        width=SLIDE_W - Inches(0.5), height=Inches(0.9),
        font_size=24, bold=True, color=C_WHITE,
        align=PP_ALIGN.LEFT,
    )


def _add_bullets(slide, bullets: list[str], top=Inches(1.15)) -> None:
    """Render bullet list below title."""
    if not bullets:
        return
    content_h = SLIDE_H - top - Inches(0.1)
    txBox = slide.shapes.add_textbox(
        Inches(0.35), top, SLIDE_W - Inches(0.5), content_h
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        # Indent sub-bullets (leading spaces or tabs)
        p.level = 0
        run = p.add_run()
        run.text = "• " + bullet
        run.font.size = Pt(18)
        run.font.color.rgb = C_BODY_TEXT
        run.font.name = "Calibri"


def _embed_images_grid(slide, images: list[Path],
                        top=Inches(1.1), bottom_margin=Inches(0.1)) -> None:
    """Tile 1-4 images in a grid within the body area."""
    if not images:
        return
    n = len(images)
    area_left = Inches(0.25)
    area_top = top
    area_w = SLIDE_W - Inches(0.5)
    area_h = SLIDE_H - top - bottom_margin

    if n == 1:
        cols, rows = 1, 1
    elif n == 2:
        cols, rows = 2, 1
    elif n <= 4:
        cols, rows = 2, 2
    else:
        cols = n
        rows = 1

    cell_w = area_w / cols
    cell_h = area_h / rows
    padding = Inches(0.05)

    for i, img_path in enumerate(images):
        col = i % cols
        row = i // cols
        img_left = area_left + col * cell_w + padding
        img_top = area_top + row * cell_h + padding
        img_w = cell_w - 2 * padding
        img_h = cell_h - 2 * padding
        try:
            slide.shapes.add_picture(
                str(img_path), img_left, img_top, img_w, img_h
            )
        except Exception as e:
            print(f"  Warning: could not embed {img_path.name}: {e}", file=sys.stderr)


def _embed_images_with_bullets(slide, bullets: list[str], images: list[Path]) -> None:
    """
    When a slide has both text and images, split the content area:
    left 45% for bullets, right 55% for images.
    """
    if not images:
        _add_bullets(slide, bullets)
        return
    if not bullets:
        _embed_images_grid(slide, images)
        return

    # Split layout
    top = Inches(1.15)
    content_h = SLIDE_H - top - Inches(0.1)
    left_w = SLIDE_W * 0.42
    right_w = SLIDE_W - left_w - Inches(0.3)
    right_left = Inches(0.15) + left_w + Inches(0.1)

    # Bullets on left
    txBox = slide.shapes.add_textbox(
        Inches(0.35), top, left_w - Inches(0.2), content_h
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = "• " + bullet
        run.font.size = Pt(16)
        run.font.color.rgb = C_BODY_TEXT
        run.font.name = "Calibri"

    # Images on right
    n = len(images)
    if n == 1:
        cols, rows = 1, 1
    elif n == 2:
        cols, rows = 1, 2
    else:
        cols, rows = 2, 2

    cell_w = right_w / cols
    cell_h = content_h / rows
    padding = Inches(0.05)

    for i, img_path in enumerate(images):
        col = i % cols
        row = i // cols
        img_left = right_left + col * cell_w + padding
        img_top = top + row * cell_h + padding
        img_w = cell_w - 2 * padding
        img_h = cell_h - 2 * padding
        try:
            slide.shapes.add_picture(str(img_path), img_left, img_top, img_w, img_h)
        except Exception as e:
            print(f"  Warning: could not embed {img_path.name}: {e}", file=sys.stderr)


# ---------------------------------------------------------------------------
# Slide rendering
# ---------------------------------------------------------------------------

def build_title_slide(prs: Presentation, spec: SlideSpec) -> None:
    """First slide: centered title layout."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Dark background
    _set_slide_bg(slide, C_DARK_BLUE)

    # Center strip for title
    strip = slide.shapes.add_shape(
        1,
        Inches(0), Inches(2.2),
        SLIDE_W, Inches(2.8),
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = C_ACCENT
    strip.fill.fore_color.theme_color  # touch to avoid lazy init issues
    strip.line.fill.background()

    # Main title
    _add_text_box(
        slide, spec.title,
        left=Inches(0.8), top=Inches(2.35),
        width=SLIDE_W - Inches(1.6), height=Inches(2.0),
        font_size=28, bold=True, color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )

    # Sub-bullets (author, supervisor, date)
    sub_text = "\n".join(spec.bullets[:3]) if spec.bullets else ""
    if sub_text:
        _add_text_box(
            slide, sub_text,
            left=Inches(1.0), top=Inches(5.3),
            width=SLIDE_W - Inches(2.0), height=Inches(1.8),
            font_size=20, bold=False, color=C_WHITE,
            align=PP_ALIGN.CENTER,
        )

    # Background image (decorative, bottom-right corner)
    if spec.images:
        img = spec.images[0]
        try:
            slide.shapes.add_picture(
                str(img),
                SLIDE_W - Inches(4.5), SLIDE_H - Inches(2.2),
                Inches(4.3), Inches(2.1),
            )
        except Exception as e:
            print(f"  Warning: title bg image {img.name}: {e}", file=sys.stderr)


def build_content_slide(prs: Presentation, spec: SlideSpec) -> None:
    """Standard content slide with title + bullets ± images."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # White background
    _set_slide_bg(slide, C_WHITE)

    # Title bar
    _add_title_box(slide, spec.title)

    # Content
    if spec.images and not spec.bullets:
        # Image-only body
        _embed_images_grid(slide, spec.images)
    elif spec.images and spec.bullets:
        # Mixed: bullets + images
        _embed_images_with_bullets(slide, spec.bullets, spec.images)
    else:
        # Text-only
        _add_bullets(slide, spec.bullets)

    # Slide number (bottom-right, small)
    num_text = f"Slide {spec.slide_id}"
    _add_text_box(
        slide, num_text,
        left=SLIDE_W - Inches(1.5), top=SLIDE_H - Inches(0.38),
        width=Inches(1.4), height=Inches(0.35),
        font_size=11, bold=False, color=RGBColor(0x99, 0x99, 0x99),
        align=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Main builder
# ---------------------------------------------------------------------------

def build_pptx(slides: list[SlideSpec]) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for i, spec in enumerate(slides):
        print(f"  Building Slide {spec.slide_id}: {spec.title[:60]}")
        if i == 0:
            build_title_slide(prs, spec)
        else:
            build_content_slide(prs, spec)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    size_kb = OUTPUT_PATH.stat().st_size // 1024
    print(f"\nSaved: {OUTPUT_PATH}  ({size_kb} KB)")


# ---------------------------------------------------------------------------
# Validate mode
# ---------------------------------------------------------------------------

def validate(slides: list[SlideSpec]) -> None:
    print(f"Total slides: {len(slides)}\n")
    for spec in slides:
        img_note = f"  [{len(spec.images)} image(s)]" if spec.images else ""
        print(f"Slide {spec.slide_id}: {spec.title}{img_note}")


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main() -> None:
    parser = argparse.ArgumentParser(description="Build thesis defense PPTX from outline.md")
    parser.add_argument("--validate", action="store_true",
                        help="Print slide list only, do not generate PPTX")
    args = parser.parse_args()

    print(f"Parsing outline: {OUTLINE_PATH}")
    slides = parse_outline(OUTLINE_PATH)
    print(f"Found {len(slides)} slides\n")

    if args.validate:
        validate(slides)
    else:
        build_pptx(slides)


if __name__ == "__main__":
    main()
